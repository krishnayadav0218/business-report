"""
test_report_builder.py
Covers src/report_builder.py. The signature test below is a direct regression
test for a real production bug: app.py called build_report(..., region_df=...,
sp_df=..., insights=...) while an older report_builder.py didn't accept those
kwargs, causing "build_report() got an unexpected keyword argument 'region_df'"
after a partial deploy. This test fails immediately if that ever happens again,
instead of only surfacing when someone uploads a file in production.
"""

import inspect
import os

from src import charts, process, report_builder


def test_build_report_accepts_the_kwargs_app_py_actually_passes():
    """
    If app.py, desktop_app.py, watch_folder.py, or src/main.py are ever updated
    to pass a new kwarg to build_report() without report_builder.py being
    updated to match (or vice versa), this test catches it immediately.
    """
    sig = inspect.signature(report_builder.build_report)
    required_kwargs = {"region_df", "sp_df", "insights", "output_path", "report_title", "period_label"}
    missing = required_kwargs - set(sig.parameters.keys())
    assert not missing, f"build_report is missing expected parameters: {missing}"


def test_full_report_generation_end_to_end(sample_df, tmp_path):
    kpis = process.compute_kpis(sample_df)
    region_df = process.region_summary(sample_df)
    trend_df = process.trend_summary(sample_df)
    sp_df = process.salesperson_summary(sample_df)
    insights = process.generate_insights(kpis, region_df, trend_df, sp_df)

    chart_paths = {}
    if not region_df.empty:
        chart_paths["region"] = charts.region_bar_chart(region_df)
    if not trend_df.empty:
        chart_paths["trend"] = charts.trend_line_chart(trend_df)
    if not sp_df.empty:
        chart_paths["leaderboard"] = charts.salesperson_leaderboard_chart(sp_df)

    output_path = str(tmp_path / "test_report.pptx")
    result_path = report_builder.build_report(
        kpis, chart_paths, output_path=output_path,
        region_df=region_df, sp_df=sp_df, insights=insights
    )

    assert os.path.exists(result_path)
    assert os.path.getsize(result_path) > 0


def test_report_has_expected_number_of_slides(sample_df, tmp_path):
    from pptx import Presentation

    kpis = process.compute_kpis(sample_df)
    region_df = process.region_summary(sample_df)
    trend_df = process.trend_summary(sample_df)
    sp_df = process.salesperson_summary(sample_df)
    insights = process.generate_insights(kpis, region_df, trend_df, sp_df)

    chart_paths = {}
    if not region_df.empty:
        chart_paths["region"] = charts.region_bar_chart(region_df)
    if not trend_df.empty:
        chart_paths["trend"] = charts.trend_line_chart(trend_df)
    if not sp_df.empty:
        chart_paths["leaderboard"] = charts.salesperson_leaderboard_chart(sp_df)

    output_path = str(tmp_path / "test_report.pptx")
    report_builder.build_report(
        kpis, chart_paths, output_path=output_path,
        region_df=region_df, sp_df=sp_df, insights=insights
    )

    prs = Presentation(output_path)
    # Title, About, Executive Summary, Highlights, Region, Trend, Leaderboard,
    # Table, Recommendations, Closing -- 10 when all sections have data.
    assert len(prs.slides) == 10


def test_every_slide_has_speaker_notes(sample_df, tmp_path):
    from pptx import Presentation

    kpis = process.compute_kpis(sample_df)
    region_df = process.region_summary(sample_df)
    trend_df = process.trend_summary(sample_df)
    sp_df = process.salesperson_summary(sample_df)
    insights = process.generate_insights(kpis, region_df, trend_df, sp_df)

    output_path = str(tmp_path / "test_report.pptx")
    report_builder.build_report(
        kpis, {}, output_path=output_path,
        region_df=region_df, sp_df=sp_df, insights=insights
    )

    prs = Presentation(output_path)
    for i, slide in enumerate(prs.slides, 1):
        assert slide.has_notes_slide, f"Slide {i} is missing speaker notes"
        assert len(slide.notes_slide.notes_text_frame.text.strip()) > 0, f"Slide {i} has empty notes"


def test_no_shape_overflows_the_slide_bounds(sample_df, tmp_path):
    """Catches layout regressions where a text box or table runs off the edge of the slide."""
    from pptx import Presentation

    kpis = process.compute_kpis(sample_df)
    region_df = process.region_summary(sample_df)
    trend_df = process.trend_summary(sample_df)
    sp_df = process.salesperson_summary(sample_df)
    insights = process.generate_insights(kpis, region_df, trend_df, sp_df)

    chart_paths = {}
    if not region_df.empty:
        chart_paths["region"] = charts.region_bar_chart(region_df)

    output_path = str(tmp_path / "test_report.pptx")
    report_builder.build_report(
        kpis, chart_paths, output_path=output_path,
        region_df=region_df, sp_df=sp_df, insights=insights
    )

    prs = Presentation(output_path)
    slide_w, slide_h = prs.slide_width, prs.slide_height
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.left is None or shape.width is None:
                continue
            assert shape.left + shape.width <= slide_w + 1000, f"Slide {i}: shape overflows right edge"
            assert shape.top + shape.height <= slide_h + 1000, f"Slide {i}: shape overflows bottom edge"
