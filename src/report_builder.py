"""
report_builder.py
Builds a meeting-ready PPTX report: title, methodology, executive summary,
key highlights, chart slides with plain-English commentary next to each
chart, a data table, and a recommendations slide.

Every slide also gets full SPEAKER NOTES (visible only to the presenter in
PowerPoint's Notes pane / Presenter View, never on the projected slide) --
a ready-to-read script so whoever presents doesn't have to improvise an
explanation of what a chart or number means.

Design note: instead of hand-editing a fixed .pptx template's placeholders
(brittle -- breaks the moment a shape gets renamed), this builds each slide
programmatically with python-pptx using fixed style helper functions (colors,
fonts, spacing). Same end result -- a consistent branded deck every run --
without depending on a specific template file surviving edits in PowerPoint.
"""

import os
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x2E, 0x9E, 0x8F)
CORAL = RGBColor(0xE4, 0x62, 0x2C)
GREY = RGBColor(0x8A, 0x94, 0xA6)
LIGHT_BG = RGBColor(0xF7, 0xF8, 0xFA)
LIGHT_LINE = RGBColor(0xE4, 0xE7, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _blank_slide(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # fully blank layout
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def _set_notes(slide, text):
    """
    Adds presenter speaker notes to a slide -- visible in PowerPoint's Notes
    pane / Presenter View, never on the slide itself when projected. This is
    what makes the deck "meeting ready": a full talking-points script for
    whoever presents, generated fresh from this run's actual numbers.
    """
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def _textbox(slide, left, top, width, height, text, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, font="Calibri", italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return box


def _bullet_list(slide, left, top, width, height, items, size=14, color=NAVY,
                  bullet_color=TEAL, line_spacing=1.35, gap_after=10):
    """A clean bullet list -- used for highlights, insights, and recommendations."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(gap_after)
        r1 = p.add_run()
        r1.text = "\u25cf  "
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.name = "Calibri"
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = "Calibri"
    return box


def _slide_header(slide, title, subtitle=None):
    _textbox(slide, MARGIN, Inches(0.45), Inches(11.5), Inches(0.6),
              title, size=26, bold=True, color=NAVY)
    if subtitle:
        _textbox(slide, MARGIN, Inches(1.02), Inches(11.5), Inches(0.4),
                  subtitle, size=13, color=GREY, italic=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.42), Inches(12.13), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_LINE
    line.line.fill.background()
    line.shadow.inherit = False


def _kpi_card(slide, left, top, width, height, label, value, caption=None, accent=TEAL):
    """A KPI card with a label, the big number, and a small plain-English
    explanation of what the metric means -- so a card is never just a number."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.fill.background()
    card.shadow.inherit = False

    _textbox(slide, left + Inches(0.25), top + Inches(0.16), width - Inches(0.5), Inches(0.3),
              label, size=13, bold=True, color=GREY)
    _textbox(slide, left + Inches(0.25), top + Inches(0.46), width - Inches(0.5), Inches(0.55),
              value, size=25, bold=True, color=accent)
    if caption:
        _textbox(slide, left + Inches(0.25), top + Inches(1.04), width - Inches(0.5), Inches(0.6),
                  caption, size=10, color=GREY, italic=True)


def _insight_panel(slide, left, top, width, height, panel_title, items):
    """A light card holding a short heading + bullet insights -- sits beside charts."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.03
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.fill.background()
    card.shadow.inherit = False

    _textbox(slide, left + Inches(0.3), top + Inches(0.22), width - Inches(0.6), Inches(0.4),
              panel_title, size=14, bold=True, color=NAVY)
    if items:
        _bullet_list(slide, left + Inches(0.3), top + Inches(0.75), width - Inches(0.6),
                     height - Inches(1.0), items, size=12.5, line_spacing=1.3, gap_after=8)
    else:
        _textbox(slide, left + Inches(0.3), top + Inches(0.75), width - Inches(0.6), Inches(0.5),
                  "No notable observations for this section.", size=12, color=GREY, italic=True)


def _footer(slide, page_label):
    _textbox(slide, MARGIN, Inches(7.12), Inches(6), Inches(0.3),
              "Confidential \u2014 for internal use", size=9, color=GREY)
    _textbox(slide, Inches(11.5), Inches(7.12), Inches(1.2), Inches(0.3),
              page_label, size=9, color=GREY, align=PP_ALIGN.RIGHT)


def _chart_with_insights_slide(prs, title, chart_path, insight_title, insights, page_label,
                                subtitle=None, notes=None):
    slide = _blank_slide(prs)
    _slide_header(slide, title, subtitle)

    chart_left = MARGIN
    chart_top = Inches(1.65)
    chart_width = Inches(7.6)
    slide.shapes.add_picture(chart_path, chart_left, chart_top, width=chart_width)

    panel_left = Inches(8.45)
    panel_width = Inches(4.28)
    _insight_panel(slide, panel_left, chart_top, panel_width, Inches(5.2), insight_title, insights)

    _footer(slide, page_label)
    if notes:
        _set_notes(slide, notes)
    return slide


def _speak(*sentences) -> str:
    """Joins sentences into a single speaker-notes paragraph."""
    return " ".join(s for s in sentences if s)


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def build_report(kpis: dict, chart_paths: dict, output_path="output/Weekly_Business_Report.pptx",
                  report_title="Business Report", period_label=None,
                  region_df=None, sp_df=None, insights=None):
    """
    kpis         -- dict from process.compute_kpis()
    chart_paths  -- dict with any of: region, trend, leaderboard (paths to PNGs)
    region_df    -- optional DataFrame from process.region_summary() (for the data table)
    sp_df        -- optional DataFrame from process.salesperson_summary() (unused directly,
                    kept for symmetry / future extension)
    insights     -- optional dict from process.generate_insights() -- if omitted, chart
                    slides fall back to full-width charts with no commentary panel
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    period_label = period_label or date.today().strftime("%d %b %Y")
    insights = insights or {}
    page_num = 1

    def next_page():
        nonlocal page_num
        page_num += 1
        return f"{page_num - 1}"

    headline = insights.get("headline")

    # ---- Slide 1: Title ----
    s1 = _blank_slide(prs, bg=NAVY)
    _textbox(s1, Inches(0.8), Inches(2.5), Inches(11), Inches(1.2),
              report_title, size=40, bold=True, color=WHITE)
    if headline:
        _textbox(s1, Inches(0.8), Inches(3.35), Inches(11), Inches(0.5),
                  headline, size=17, color=TEAL, bold=True)
    _textbox(s1, Inches(0.8), Inches(3.95), Inches(11), Inches(0.5),
              f"Reporting period ending {period_label}", size=14, color=RGBColor(0xB8, 0xC2, 0xD6))
    _set_notes(s1, _speak(
        "Welcome everyone to this business review.",
        f"This report covers the period ending {period_label}." if period_label else "",
        f"The one-line summary going in: {headline}." if headline else "",
        "We'll walk through the numbers, what's driving them, and what we should do next."
    ))

    # ---- Slide 2: About This Report (methodology) ----
    s2 = _blank_slide(prs)
    _slide_header(s2, "About This Report", "How the figures on the following slides were calculated")
    methodology_points = [
        "Generated directly from the uploaded file \u2014 no manual entry or adjustment to the numbers.",
        "Target Achievement = Total Collection \u00f7 Total Target \u00d7 100.",
        "Profit / Loss = Total Collection \u2212 Total Cost.",
        "Average Order Value = Total Collection \u00f7 Total Orders.",
        "Period-on-Period Growth compares the first half of the reporting period to the second half.",
        "All amounts are shown in the currency used in the source file (\u20b9 by default).",
    ]
    _bullet_list(s2, MARGIN, Inches(1.9), Inches(12.1), Inches(4.6),
                 methodology_points, size=15, line_spacing=1.5, gap_after=14)
    _footer(s2, next_page())
    _set_notes(s2, _speak(
        "Before the numbers, a quick word on where they come from.",
        "Everything in this deck is calculated straight from the uploaded file \u2014 nothing has been manually adjusted.",
        "If anyone asks how a specific figure was derived, the formulas are on this slide:",
        "achievement is collection over target, profit is collection minus cost,",
        "average order value is collection over the number of orders,",
        "and the period-on-period growth compares the first half of the period to the second half.",
        "Keep this slide in mind \u2014 we'll refer back to these definitions as we go."
    ))

    # ---- Slide 3: Executive Summary (KPI cards) ----
    s3 = _blank_slide(prs)
    _slide_header(s3, "Executive Summary", headline)

    card_w, card_h = Inches(2.9), Inches(1.75)
    gap = Inches(0.3)
    top = Inches(1.75)
    cards = [
        ("Total Collection", f"\u20b9{kpis['total_collection']:,.0f}",
         "Total amount collected this period", TEAL),
        ("Target Achievement", f"{kpis['achievement_pct']:.1f}%",
         "Share of the target reached", NAVY),
        ("Total Orders", f"{kpis['total_orders']:,}",
         "Transactions processed this period", CORAL),
        ("Avg Order Value", f"\u20b9{kpis['aov']:,.0f}",
         "Average revenue per order", TEAL),
    ]
    left = MARGIN
    for label, value, caption, accent in cards:
        _kpi_card(s3, left, top, card_w, card_h, label, value, caption, accent)
        left = Emu(left + card_w + gap)

    profit_color = TEAL if kpis["profit"] >= 0 else CORAL
    profit_label = "Profit" if kpis["profit"] >= 0 else "Loss"
    _kpi_card(s3, MARGIN, Inches(3.75), card_w, card_h, profit_label,
              f"\u20b9{abs(kpis['profit']):,.0f}", "Total collection minus total cost", profit_color)

    if kpis.get("wow_growth") is not None:
        growth_color = TEAL if kpis["wow_growth"] >= 0 else CORAL
        arrow = "\u25b2" if kpis["wow_growth"] >= 0 else "\u25bc"
        _kpi_card(s3, Emu(MARGIN + card_w + gap), Inches(3.75), card_w, card_h,
                  "Period-on-Period Growth", f"{arrow} {abs(kpis['wow_growth']):.1f}%",
                  "Change vs. the earlier half of this period", growth_color)
    _footer(s3, next_page())

    growth_note = ""
    if kpis.get("wow_growth") is not None:
        direction = "up" if kpis["wow_growth"] >= 0 else "down"
        growth_note = f"Period-on-period, collection is {direction} {abs(kpis['wow_growth']):.1f} percent."
    _set_notes(s3, _speak(
        f"Here's the top-line picture. Total collection came in at \u20b9{kpis['total_collection']:,.0f}",
        f"against a target of \u20b9{kpis['total_target']:,.0f}, which is {kpis['achievement_pct']:.1f} percent of target." if kpis.get('total_target') else "",
        f"We processed {kpis['total_orders']:,} orders, averaging \u20b9{kpis['aov']:,.0f} per order." if kpis.get('total_orders') else "",
        f"On profitability: {'a net profit' if kpis['profit'] >= 0 else 'a net loss'} of \u20b9{abs(kpis['profit']):,.0f}.",
        growth_note,
        "Pause here for questions before moving into the detail."
    ))

    # ---- Slide 4: Key Highlights ----
    if insights.get("highlights"):
        s4 = _blank_slide(prs)
        _slide_header(s4, "Key Highlights", "What stands out in this period's numbers")
        _bullet_list(s4, MARGIN, Inches(1.9), Inches(12.1), Inches(4.8),
                     insights["highlights"], size=16, line_spacing=1.5, gap_after=16)
        _footer(s4, next_page())
        _set_notes(s4, _speak(
            "These are the points worth calling out explicitly \u2014 read through each one:",
            " ".join(f"({i+1}) {h}" for i, h in enumerate(insights["highlights"])),
            "Any of these worth a deeper conversation before we move to the regional and trend detail?"
        ))

    # ---- Slide 5: Regional Performance ----
    if chart_paths.get("region"):
        region_notes = _speak(
            "This chart compares target against actual collection for each region.",
            " ".join(insights.get("region", [])),
            "Use this to steer where follow-up attention or extra support should go."
        )
        _chart_with_insights_slide(
            prs, "Regional Performance", chart_paths["region"],
            "What this shows", insights.get("region", []), next_page(),
            subtitle="Target vs. collection across regions", notes=region_notes
        )

    # ---- Slide 6: Collection Trend ----
    if chart_paths.get("trend"):
        trend_notes = _speak(
            "This line tracks daily collection across the reporting period.",
            " ".join(insights.get("trend", [])),
            "Flag any specific day here if there's a known reason behind a spike or dip (e.g. a promotion, holiday, or system issue)."
        )
        _chart_with_insights_slide(
            prs, "Collection Trend", chart_paths["trend"],
            "What this shows", insights.get("trend", []), next_page(),
            subtitle="Daily collection over the reporting period", notes=trend_notes
        )

    # ---- Slide 7: Salesperson Leaderboard ----
    if chart_paths.get("leaderboard"):
        sp_notes = _speak(
            "This ranks individual contribution to total collection.",
            " ".join(insights.get("salesperson", [])),
            "Good moment to publicly recognize the top performer, and to check in privately with anyone trailing the team average."
        )
        _chart_with_insights_slide(
            prs, "Salesperson Leaderboard", chart_paths["leaderboard"],
            "What this shows", insights.get("salesperson", []), next_page(),
            subtitle="Individual contribution to total collection", notes=sp_notes
        )

    # ---- Slide 8: Detailed Breakdown Table ----
    if region_df is not None and not region_df.empty:
        s8 = _blank_slide(prs)
        _slide_header(s8, "Regional Breakdown", "Detailed figures by region")

        has_target = "Target" in region_df.columns
        headers = ["Region", "Collection"] + (["Target", "Achievement"] if has_target else [])
        rows = len(region_df) + 1
        cols = len(headers)
        table_shape = s8.shapes.add_table(rows, cols, MARGIN, Inches(1.75), Inches(12.1), Inches(0.5) * rows)
        table = table_shape.table

        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(13)
            para.font.bold = True
            para.font.color.rgb = WHITE

        below_target = []
        for r, (_, row) in enumerate(region_df.iterrows(), start=1):
            values = [str(row["Region"]), f"\u20b9{row['Collection']:,.0f}"]
            if has_target:
                achievement = (row["Collection"] / row["Target"] * 100) if row["Target"] else 0
                values += [f"\u20b9{row['Target']:,.0f}", f"{achievement:.1f}%"]
                if achievement < 100:
                    below_target.append(str(row["Region"]))
            for c, val in enumerate(values):
                cell = table.cell(r, c)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG if r % 2 == 0 else WHITE
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(12)
                para.font.color.rgb = NAVY
        _footer(s8, next_page())
        _set_notes(s8, _speak(
            "This is the detailed region-by-region breakdown \u2014 use it if someone asks for specifics beyond the summary chart.",
            f"Regions below their target this period: {', '.join(below_target)}." if below_target else "Every region met or exceeded its target this period.",
        ))

    # ---- Slide 9: Recommendations ----
    if insights.get("recommendations"):
        s9 = _blank_slide(prs)
        _slide_header(s9, "Recommendations", "Suggested next steps based on this period's data")
        _bullet_list(s9, MARGIN, Inches(1.9), Inches(12.1), Inches(4.8),
                     insights["recommendations"], size=16, line_spacing=1.6, gap_after=18,
                     bullet_color=CORAL)
        _footer(s9, next_page())
        _set_notes(s9, _speak(
            "These recommendations follow directly from what we just walked through.",
            " ".join(f"({i+1}) {r}" for i, r in enumerate(insights["recommendations"])),
            "Where possible, assign an owner and a follow-up date to each item before closing the meeting."
        ))

    # ---- Slide 10: Closing ----
    s10 = _blank_slide(prs, bg=NAVY)
    _textbox(s10, Inches(0.8), Inches(3.2), Inches(11), Inches(1),
              "Thank You", size=36, bold=True, color=WHITE)
    _textbox(s10, Inches(0.8), Inches(4.0), Inches(11), Inches(0.5),
              "Questions and discussion welcome.", size=15, color=RGBColor(0xB8, 0xC2, 0xD6))
    _set_notes(s10, _speak(
        "That's the full picture for this period.",
        "Open the floor for questions, confirm owners for each recommendation, and agree on when we'll next review these numbers."
    ))

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path
